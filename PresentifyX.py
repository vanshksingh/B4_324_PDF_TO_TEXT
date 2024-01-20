import sys
import os
import shutil
import cv2
import numpy as np
from PIL import Image
import keras_ocr
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt, Emu
from pdf2docx import Converter
from transformers import T5Tokenizer, T5ForConditionalGeneration

from PyQt5.QtCore import Qt, QSize
from PyQt5.QtGui import QPixmap, QImage, QFont, QColor, QPalette
from PyQt5.QtWidgets import (QApplication, QWidget, QVBoxLayout, QHBoxLayout, QPushButton, QLabel, QFileDialog, QTextEdit, QComboBox)

class Summarization:
    @staticmethod
    def generate_summary(text, max_length=150, model_name="t5-base"):
        # Initialize the tokenizer and model
        tokenizer = T5Tokenizer.from_pretrained(model_name)
        model = T5ForConditionalGeneration.from_pretrained(model_name)

        # Preprocess and encode the text for summarization
        inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=512, truncation=True)

        # Generate summary ids
        summary_ids = model.generate(inputs, max_length=max_length, length_penalty=2.0, num_beams=4,
                                     early_stopping=True)

        # Decode and return the summary
        summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        return summary

    @staticmethod
    def summarize_text(text_to_summarize):
        return Summarization.generate_summary(text_to_summarize)


def process_input(input_path):
    if input_path.lower().endswith('.pdf'):
        all_text, output_path = pdf_to_ppt(input_path)
        return all_text, output_path
    else:
        all_text, output_path = image_to_ppt(input_path)
        return all_text, output_path


def pdf_to_images(pdf_path, output_folder='pdf_images'):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    doc = fitz.open(pdf_path)
    image_paths = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        image_path = os.path.join(output_folder, f'page_{page_num}.png')
        pix.save(image_path)
        image_paths.append(image_path)

    doc.close()
    return image_paths


def pdf_to_ppt(pdf_path, output_ppt_path=None):
    if output_ppt_path is None:
        output_ppt_path = os.path.splitext(pdf_path)[0] + '_presentation.pptx'

    image_paths = pdf_to_images(pdf_path)
    prs = Presentation()

    all_detected_text = []

    for image_path in image_paths:
        detected_text = remove_text_and_create_ppt(image_path, prs)
        all_detected_text.append(detected_text)

    prs.save(output_ppt_path)

    return ' '.join(all_detected_text), output_ppt_path


def image_to_ppt(image_path, output_ppt_path=None):
    if output_ppt_path is None:
        output_ppt_path = os.path.splitext(image_path)[0] + '_presentation.pptx'

    prs = Presentation()
    detected_text = remove_text_and_create_ppt(image_path, prs)
    prs.save(output_ppt_path)

    return detected_text, output_ppt_path


def pixels_to_emus(pixel, dpi=96):
    inches = pixel / dpi
    return Emu(inches * 914400)


def estimate_font_size(height, scaling_factor=0.75):
    return Pt(height * scaling_factor)


def midpoint(x1, y1, x2, y2):
    return (int((x1 + x2) / 2), int((y1 + y2) / 2))


def remove_text_and_create_ppt(image_path, prs):
    pipeline = keras_ocr.pipeline.Pipeline()

    image = keras_ocr.tools.read(image_path)
    h, w, _ = image.shape

    prediction_groups = pipeline.recognize([image])

    mask = np.zeros(image.shape[:2], dtype="uint8")
    text_positions = []
    all_detected_text = []
    for box in prediction_groups[0]:
        text, box_points = box
        all_detected_text.append(text)

        x0, y0, x1, y1, x2, y2, x3, y3 = np.array(box_points).flatten()

        x_mid0, y_mid0 = midpoint(x1, y1, x2, y2)
        x_mid1, y_mid1 = midpoint(x0, y0, x3, y3)

        thickness = int(np.sqrt((x2 - x1) ** 2 + (y2 - y1) ** 2))

        cv2.line(mask, (x_mid0, y_mid0), (x_mid1, y_mid1), 255, thickness)
        text_positions.append((text, x0, y0, x2 - x0, y2 - y0))

    result = cv2.inpaint(image, mask, 7, cv2.INPAINT_NS)
    output_image_path = os.path.splitext(image_path)[0] + '_output.jpg'
    cv2.imwrite(output_image_path, cv2.cvtColor(result, cv2.COLOR_BGR2RGB))

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(output_image_path, Emu(0), Emu(0), width=pixels_to_emus(w), height=pixels_to_emus(h))

    for text, x, y, width, height in text_positions:
        textbox = slide.shapes.add_textbox(pixels_to_emus(x), pixels_to_emus(y),
                                           pixels_to_emus(width), pixels_to_emus(height))
        p = textbox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = estimate_font_size(height)

    return ' '.join(all_detected_text)


class SampleApp(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()
        self.selectedFilePath = ""

    def initUI(self):
        # Color scheme settings
        primary_color = QColor(10, 25, 47, 210)  # Dark blue with some transparency for a glass-like effect
        button_color = "rgba(0, 121, 191, 180)"  # Medium dark blue for buttons
        border_color = "#0079BF"  # Medium dark blue, matching the buttons
        text_area_color = "#17202A"  # Very dark blue, almost black
        preview_box_color = "black"  # Pure black for the preview box
        preview_box_border_color = "#0079BF"  # Medium dark blue to match buttons

        # Main vertical layout
        main_layout = QVBoxLayout(self)

        # Horizontal layout for the top part (buttons and preview)
        top_layout = QHBoxLayout()

        # Left side layout for buttons
        left_layout = QVBoxLayout()

        # Standard font for buttons
        button_font = QFont('Arial', 12)

        # Add title label with custom styling
        title_label = QLabel('PresentifyX')
        title_font = QFont('Courier New', 30, QFont.Bold)
        title_font.setUnderline(True)
        title_label.setFont(title_font)
        title_label.setStyleSheet("color: red;")
        title_label.setAlignment(Qt.AlignCenter)
        left_layout.addWidget(title_label, alignment=Qt.AlignCenter)

        # Spacer below title
        left_layout.addSpacing(40)

        # Buttons and other controls
        self.summarize_button = QPushButton('Summarize Text')
        self.summarize_button.setFont(button_font)
        self.summarize_button.clicked.connect(self.summarize_text)
        self.summarize_button.setStyleSheet(f"""
            QPushButton {{
                color: white;
                background-color: {button_color};
                border: 2px solid {border_color};
                border-radius: 5px;
            }}
        """)
        left_layout.addWidget(self.summarize_button)

        self.add_file_button = QPushButton('Add File')
        self.add_file_button.setFont(button_font)
        self.add_file_button.clicked.connect(self.add_file)
        self.add_file_button.setStyleSheet(f"""
            QPushButton {{
                color: white;
                background-color: {button_color};
                border: 2px solid {border_color};
                border-radius: 5px;
            }}
        """)
        left_layout.addWidget(self.add_file_button)

        self.drop_down_menu = QComboBox()
        self.drop_down_menu.addItems(["PDF/IMG to PPT", "PDF to WORD", "IMG to PDF"])
        self.drop_down_menu.setStyleSheet(f"""
            QComboBox {{
                color: white;
                background-color: {button_color};
                border: 2px solid {border_color};
                border-radius: 5px;
            }}
            QComboBox QAbstractItemView {{
                color: white;
                selection-background-color: darkgray;
            }}
        """)
        left_layout.addWidget(self.drop_down_menu)

        self.start_button = QPushButton('Start')
        self.start_button.setFont(button_font)
        self.start_button.clicked.connect(self.start_action)
        self.start_button.setStyleSheet(f"""
            QPushButton {{
                color: white;
                background-color: {button_color};
                border: 2px solid {border_color};
                border-radius: 5px;
            }}
        """)
        left_layout.addWidget(self.start_button)

        self.status_label = QLabel('Status: Idle')
        left_layout.addWidget(self.status_label)

        # Add left layout to the top layout with a stretch factor
        top_layout.addLayout(left_layout, 1)

        # Right side for PDF preview
        self.preview_label = QLabel('Preview\n(Placeholder)')
        self.preview_label.setAlignment(Qt.AlignCenter)
        self.preview_label.setStyleSheet(f"""
            border: 1px solid {preview_box_border_color};
            border-radius: 10px;  
            background-color: {preview_box_color};
        """)

        top_layout.addWidget(self.preview_label, 1)

        # Add top layout to the main layout
        main_layout.addLayout(top_layout)

        # Bottom layout for the text area
        bottom_layout = QVBoxLayout()
        self.text_area = QTextEdit()
        self.text_area.setFont(QFont("Arial", 14))
        self.text_area.setStyleSheet(f"""
            border: 1px solid {preview_box_border_color};
            border-radius: 10px;  
            color: orange;
            background-color: {preview_box_color};
        """)
        bottom_layout.addWidget(self.text_area)

        # Add bottom layout to the main layout with a stretch factor
        main_layout.addLayout(bottom_layout, 1)

        # Set the main layout
        self.setLayout(main_layout)
        self.setWindowTitle('B4/324')

        # Glass-like background effect with custom color
        self.setWindowOpacity(0.9)
        self.setAttribute(Qt.WA_TranslucentBackground, True)
        palette = QPalette()
        palette.setColor(QPalette.Window, primary_color)
        self.setPalette(palette)

        # Set a fixed window size
        self.setGeometry(100, 100, 800, 600)


    def img_to_pdf(img_path, pdf_path=None):
        if pdf_path is None:
            pdf_path = os.path.splitext(img_path)[0] + '.pdf'

        image = Image.open(img_path)
        pdf_image = image.convert('RGB')
        pdf_image.save(pdf_path)

        return pdf_path

    def pdf_to_word(self, pdf_path, word_path=None):
        if word_path is None:
            word_path = os.path.splitext(pdf_path)[0] + '.docx'

        cv = Converter(pdf_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()

        return word_path

    def summarize_text(self):
        input_text = self.text_area.toPlainText()
        if input_text.strip():
            summary = Summarization.summarize_text(input_text)
            self.text_area.setText(summary)
        else:
            self.text_area.setText("Please enter some text to summarize.")

    def download_action(self):
        if hasattr(self, 'processed_file_path') and self.processed_file_path:
            # Assuming processed_file_path is a full path to the processed file
            save_path = os.path.join(os.path.dirname(self.selectedFilePath), os.path.basename(self.processed_file_path))

            try:
                # If the processed file is different from the original file, copy it to the original file's directory
                if self.processed_file_path != save_path:
                    shutil.copy2(self.processed_file_path, save_path)

                self.update_status(f"File saved as {save_path}")
                print(f"File saved as {save_path}")
            except Exception as e:
                self.update_status("Error saving file.")
                print(f"Error: {e}")
        else:
            self.update_status("No processed file to save.")
            print("No processed file to save.")

    def add_file(self):
        file_filter = "All Supported Files (*.png *.jpg *.jpeg *.bmp *.gif *.tiff *.xpm *.pdf *.pptx *.ppt *.docx *.doc);;Image Files (*.png *.jpg *.jpeg *.bmp *.gif *.tiff *.xpm);;PDF Files (*.pdf);;PowerPoint Files (*.pptx *.ppt);;Word Files (*.docx *.doc)"
        filename, _ = QFileDialog.getOpenFileName(self, "Open File", "", file_filter)
        if filename:
            self.selectedFilePath = filename
            self.text_area.append(f"Selected File: {filename}")
            self.update_preview(filename)

    def start_action(self):

        if self.selectedFilePath:
            self.update_status("Running...")

            operation = self.drop_down_menu.currentText()
            if operation == "PDF/IMG to PPT":
                self.update_status("Running...")
                self.process_file_for_ppt()
            elif operation == "PPT to PDF":
                self.update_status("Running...")
                output_file_path = self.ppt_to_pdf(self.selectedFilePath)
                self.processed_file_path = output_file_path
                self.text_area.append(f"Converted to PDF: {output_file_path}")
            elif operation == "IMG to PDF":
                self.update_status("Running...")
                output_file_path = self.img_to_pdf(self.selectedFilePath)
                self.processed_file_path = output_file_path
                self.text_area.append(f"Converted to PDF: {output_file_path}")
            # ... existing code
            elif operation == "PDF to WORD":  # Check for PDF to Word operation
                self.update_status("Running...")
                output_file_path = self.pdf_to_word(self.selectedFilePath)
                self.processed_file_path = output_file_path
                self.text_area.append(f"Converted to Word: {output_file_path}")
            # Add more conditions here for other operations
            self.update_status("Finished")
        else:
            self.update_status("No file selected.")

    def process_file_for_ppt(self):
        if self.selectedFilePath.lower().endswith(('.pdf', '.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.xpm')):
            all_text_from_input, output_file_path = process_input(self.selectedFilePath)
            self.text_area.append("Detected Text: " + all_text_from_input)
            self.processed_file_path = output_file_path  # Save the output file path
        else:
            self.text_area.append("File format not supported for this operation.")

    def update_status(self, value):
        self.status_label.setText(f"Status: {value}")

    def update_preview(self, filename):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.xpm')):
            pixmap = QPixmap(filename)
            self.preview_label.setPixmap(pixmap.scaled(300, 300, Qt.KeepAspectRatio))  # Larger preview

        elif filename.lower().endswith('.pdf'):
            self.preview_pdf(filename)
        elif filename.lower().endswith(('.pptx', '.ppt')):
            self.preview_label.setText('PowerPoint File Selected')
        elif filename.lower().endswith(('.docx', '.doc')):
            self.preview_label.setText('Word Document Selected')
        else:
            self.preview_label.clear()

    def preview_pdf(self, filename):
        doc = fitz.open(filename)
        page = doc.load_page(0)  # first page
        pix = page.get_pixmap()
        img = QImage(pix.samples, pix.width, pix.height, pix.stride, QImage.Format_RGB888)
        pixmap = QPixmap.fromImage(img)
        self.preview_label.setPixmap(pixmap.scaled(100, 100, Qt.KeepAspectRatio))
        doc.close()


def main():
    app = QApplication(sys.argv)
    app.setApplicationName('PresentifyX')
    app.setStyle('Fusion')
    ex = SampleApp()
    ex.show()
    sys.exit(app.exec_())


if __name__ == '__main__':
    main()
